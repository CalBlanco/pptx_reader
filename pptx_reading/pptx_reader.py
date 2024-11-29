from pptx import Presentation
import os

def get_text_from_pptx(pptx_file:str)->list[tuple[str, int, str]]:
    """Load in text based items from pptx
    
    *Args*:
        pptx_file<str>: file path to find pptx at
    
    *Returns*
        A list of tuples where each tuple has the shape: (file_name, page number, text from page)

    *Notes*: Only utilizes actual text (text from images to come soon)
    """

    assert pptx_file is not None

    
    file_name = os.path.basename(pptx_file) # get the file name itself

    prs = Presentation(pptx_file) #turn into loadable pptx

    slide_data = [] 
    for page, slide in enumerate(prs.slides): #Collect page text
        page_text = []
        for shape in slide.shapes: #Iterate over all boxes in page

            if shape.has_text_frame: #if text box, take its text
                t = shape.text_frame
                try:
                    page_text.append(t.text)
                except:
                    continue
        slide_data.append((file_name, page+1, ' '.join(page_text))) #page is incremented by 1 to represent accurate slide in pptx

    return slide_data
